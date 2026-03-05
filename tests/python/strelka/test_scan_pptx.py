"""Tests for scan_pptx scanner."""
import io
import struct
import unittest
import zipfile
from unittest.mock import MagicMock, patch

from strelka.scanners.scan_pptx import (
    _lookup_clsid,
    _parse_ppaction,
    _extract_ole_native_info,
    _normalize_pptx_bytes,
)
import strelka.scanners.scan_pptx as scan_pptx_module


def _build_ole10native_stream(filename, src_path, temp_path, payload=b"\x4d\x5a"):
    """Build a minimal \x01Ole10Native binary stream for testing."""
    return (
        struct.pack("<I", len(payload) + 100)        # native_data_size
        + struct.pack("<H", 0x0002)                  # unknown_short
        + filename.encode() + b"\x00"                # filename (null-terminated)
        + src_path.encode() + b"\x00"                # src_path (null-terminated)
        + struct.pack("<II", 0, 0)                   # unknown_long_1, unknown_long_2
        + temp_path.encode() + b"\x00"               # temp_path (null-terminated)
        + struct.pack("<I", len(payload))             # actual_size
        + payload
    )


def _mock_olefile(stream_bytes):
    """Return a context-manager mock of olefile.OleFileIO with one Ole10Native stream."""
    mock_ole = MagicMock()
    mock_ole.__enter__ = lambda s: s
    mock_ole.__exit__ = MagicMock(return_value=False)
    mock_ole.exists.return_value = True
    mock_stream = MagicMock()
    mock_stream.read.return_value = stream_bytes
    mock_ole.openstream.return_value = mock_stream
    return mock_ole


class TestParsePpaction(unittest.TestCase):
    """Tests for _parse_ppaction helper."""

    def test_program_verb_uses_params_key(self):
        """_parse_ppaction returns 'params' not 'fields'."""
        result = _parse_ppaction("ppaction://program?file=cmd.exe")
        self.assertIn("params", result)
        self.assertNotIn("fields", result)

    def test_program_verb_and_params_parsed(self):
        """_parse_ppaction extracts verb and query params."""
        result = _parse_ppaction("ppaction://program?file=cmd.exe")
        self.assertEqual(result["verb"], "program")
        self.assertEqual(result["params"]["file"], "cmd.exe")

    def test_macro_verb(self):
        """_parse_ppaction extracts macro verb and name param."""
        result = _parse_ppaction("ppaction://macro?name=Module1.Auto_Open")
        self.assertEqual(result["verb"], "macro")
        self.assertEqual(result["params"]["name"], "Module1.Auto_Open")

    def test_ole_verb_no_params(self):
        """_parse_ppaction handles ppaction://ole with no params."""
        result = _parse_ppaction("ppaction://ole")
        self.assertEqual(result["verb"], "ole")
        self.assertEqual(result["params"], {})

    def test_non_ppaction_returns_empty(self):
        """_parse_ppaction returns null ppaction_url for non-ppaction string."""
        result = _parse_ppaction("http://example.com")
        self.assertIsNone(result["ppaction_url"])
        self.assertEqual(result["params"], {})
        self.assertIsNone(result["verb"])

    def test_none_input(self):
        """_parse_ppaction handles None input."""
        result = _parse_ppaction(None)
        self.assertIsNone(result["ppaction_url"])
        self.assertEqual(result["params"], {})


class TestLookupClsid(unittest.TestCase):
    """Tests for _lookup_clsid helper."""

    def test_known_clsid_returns_string(self):
        """_lookup_clsid returns a plain string for known CLSIDs."""
        result = _lookup_clsid("00021700-0000-0000-C000-000000000046")
        self.assertIsInstance(result, str)
        self.assertIn("Equation", result)

    def test_returns_string_not_tuple(self):
        """_lookup_clsid must not return a tuple."""
        result = _lookup_clsid("00020810-0000-0000-C000-000000000046")
        self.assertNotIsInstance(result, tuple)

    def test_unknown_clsid_returns_none(self):
        """_lookup_clsid returns None for unrecognized CLSIDs."""
        result = _lookup_clsid("DEADBEEF-0000-0000-0000-000000000000")
        self.assertIsNone(result)

    def test_braced_clsid_normalized(self):
        """_lookup_clsid handles CLSIDs wrapped in braces."""
        result = _lookup_clsid("{00021700-0000-0000-C000-000000000046}")
        self.assertIsNotNone(result)
        self.assertIn("Equation", result)


class TestHighRiskRemoved(unittest.TestCase):
    """HIGH_RISK_PROGIDS and is_high_risk opinion logic must be removed."""

    def test_high_risk_progids_not_in_module(self):
        """HIGH_RISK_PROGIDS must not exist in the module."""
        self.assertFalse(hasattr(scan_pptx_module, "HIGH_RISK_PROGIDS"))

    def test_lookup_clsid_has_no_is_dangerous(self):
        """_lookup_clsid must not return is_dangerous in any form."""
        result = _lookup_clsid("00021700-0000-0000-C000-000000000046")
        # Result must be a plain string (or None), not a tuple containing a bool
        self.assertNotIsInstance(result, tuple)


class TestScanPptxActionType(unittest.TestCase):
    """Action dicts must include action_type and use params not fields."""

    def _make_scanner(self):
        from strelka.scanners.scan_pptx import ScanPptx
        scanner = ScanPptx.__new__(ScanPptx)
        scanner.event = {}
        scanner.flags = []
        scanner.files = []
        scanner.coordinator = MagicMock()
        scanner.name = "ScanPptx"
        scanner.scanner_timeout = 30
        return scanner

    def _make_pptx_with_shape_click(self, url="http://example.com"):
        """Build a minimal PPTX with a shape-level click action hyperlink."""
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape.click_action.hyperlink.address = url
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def test_actions_have_action_type(self):
        """Every action dict must contain an 'action_type' key."""
        scanner = self._make_scanner()
        data = self._make_pptx_with_shape_click("http://evil.com")
        scanner.scan(data, MagicMock(), {}, 9999999999)
        actions = scanner.event.get("actions", [])
        self.assertTrue(len(actions) > 0, "Expected at least one action to be extracted")
        for action in actions:
            self.assertIn("action_type", action, f"action missing action_type: {action}")

    def test_actions_use_params_not_fields(self):
        """Action dicts must use 'params' key, not 'fields'."""
        scanner = self._make_scanner()
        data = self._make_pptx_with_shape_click("http://evil.com")
        scanner.scan(data, MagicMock(), {}, 9999999999)
        for action in scanner.event.get("actions", []):
            self.assertNotIn("fields", action, f"action still uses 'fields': {action}")
            self.assertIn("params", action)

    def test_hyperlink_action_type_is_hyperlink(self):
        """A plain click hyperlink action must have action_type='hyperlink'."""
        scanner = self._make_scanner()
        data = self._make_pptx_with_shape_click("http://evil.com")
        scanner.scan(data, MagicMock(), {}, 9999999999)
        actions = scanner.event.get("actions", [])
        self.assertTrue(len(actions) > 0)
        # Shape click hyperlinks with no ppaction URL → action_type=hyperlink
        hyperlink_actions = [a for a in actions if a.get("action_type") == "hyperlink"]
        self.assertTrue(len(hyperlink_actions) > 0)

    def test_event_has_no_high_risk_ole_flag(self):
        """Scanner event must not contain has_high_risk_ole."""
        scanner = self._make_scanner()
        data = self._make_pptx_with_shape_click()
        scanner.scan(data, MagicMock(), {}, 9999999999)
        self.assertNotIn("has_high_risk_ole", scanner.event)


class TestNormalizePptxBytes(unittest.TestCase):
    """_normalize_pptx_bytes must coerce slideshow/template content types."""

    PPTX_CT = "application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"

    def _make_zip_with_content_type(self, content_type):
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            ct_xml = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                f'<Override PartName="/ppt/presentation.xml" ContentType="{content_type}"/>'
                "</Types>"
            )
            zf.writestr("[Content_Types].xml", ct_xml)
        return buf.getvalue()

    def _get_content_type(self, data):
        import zipfile as _zf
        with _zf.ZipFile(io.BytesIO(data)) as zf:
            ct_xml = zf.read("[Content_Types].xml").decode()
        import re
        m = re.search(r'ContentType="([^"]+)"', ct_xml)
        return m.group(1) if m else None

    def test_ppsm_coerced_to_pptx(self):
        """Macro-enabled slideshow (.ppsm) content type is replaced with pptx."""
        data = self._make_zip_with_content_type(
            "application/vnd.ms-powerpoint.slideshow.macroEnabled.main+xml"
        )
        result, _ = _normalize_pptx_bytes(data)
        self.assertEqual(self._get_content_type(result), self.PPTX_CT)

    def test_ppsx_coerced_to_pptx(self):
        """Slideshow (.ppsx) content type is replaced with pptx."""
        data = self._make_zip_with_content_type(
            "application/vnd.openxmlformats-officedocument.presentationml.slideshow.main+xml"
        )
        result, _ = _normalize_pptx_bytes(data)
        self.assertEqual(self._get_content_type(result), self.PPTX_CT)

    def test_potx_coerced_to_pptx(self):
        """Template (.potx) content type is replaced with pptx."""
        data = self._make_zip_with_content_type(
            "application/vnd.openxmlformats-officedocument.presentationml.template.main+xml"
        )
        result, _ = _normalize_pptx_bytes(data)
        self.assertEqual(self._get_content_type(result), self.PPTX_CT)

    def test_pptx_unchanged(self):
        """Standard pptx content type is returned as-is."""
        data = self._make_zip_with_content_type(self.PPTX_CT)
        result, removed = _normalize_pptx_bytes(data)
        self.assertEqual(self._get_content_type(result), self.PPTX_CT)
        self.assertEqual(removed, 0)

    def test_pptm_unchanged(self):
        """Macro-enabled pptm content type is left as-is (python-pptx accepts it)."""
        pptm_ct = "application/vnd.ms-powerpoint.presentation.macroEnabled.main+xml"
        data = self._make_zip_with_content_type(pptm_ct)
        result, _ = _normalize_pptx_bytes(data)
        self.assertEqual(self._get_content_type(result), pptm_ct)

    def test_malformed_rel_missing_target_is_stripped(self):
        """<Relationship> with no Target attribute is removed and counted."""
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("[Content_Types].xml",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                f'<Override PartName="/ppt/presentation.xml" ContentType="{self.PPTX_CT}"/>'
                "</Types>")
            zf.writestr("ppt/slides/_rels/slide1.xml.rels",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
                '<Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink" TargetMode="External"/>'
                "</Relationships>")
        data = buf.getvalue()
        result, removed = _normalize_pptx_bytes(data)
        self.assertEqual(removed, 1)
        with zipfile.ZipFile(io.BytesIO(result)) as zf:
            rels_xml = zf.read("ppt/slides/_rels/slide1.xml.rels").decode()
        self.assertNotIn('Id="rId2"', rels_xml)
        self.assertIn('Id="rId1"', rels_xml)

    def test_malformed_rel_returns_zero_when_none(self):
        """Returns removed=0 when all relationships are well-formed."""
        buf = io.BytesIO()
        with zipfile.ZipFile(buf, "w") as zf:
            zf.writestr("[Content_Types].xml",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
                f'<Override PartName="/ppt/presentation.xml" ContentType="{self.PPTX_CT}"/>'
                "</Types>")
            zf.writestr("ppt/slides/_rels/slide1.xml.rels",
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
                '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>'
                "</Relationships>")
        data = buf.getvalue()
        _, removed = _normalize_pptx_bytes(data)
        self.assertEqual(removed, 0)

    def test_scanner_flags_malformed_relationships(self):
        """Scanner sets malformed_relationships_N flag when rels are stripped."""
        from strelka.scanners.scan_pptx import ScanPptx
        path = "/home/coz/Downloads/pptxsamples/noRelationship3.pptx"
        import os
        if not os.path.exists(path):
            self.skipTest("sample file not available")
        scanner = ScanPptx.__new__(ScanPptx)
        scanner.event = {}
        scanner.flags = []
        scanner.files = []
        scanner.coordinator = MagicMock()
        scanner.name = "ScanPptx"
        scanner.scanner_timeout = 30
        scanner.scan(open(path, "rb").read(), MagicMock(), {}, 9999999999)
        self.assertTrue(any("malformed_relationships" in f for f in scanner.flags))
        self.assertNotIn("processing_error", scanner.flags)

    def test_convocation_ppsm_loads_after_normalization(self):
        """CONVOCATION.ppsm__.pptx must load without ValueError after normalization."""
        import os
        path = "/home/coz/Downloads/pptxsamples/CONVOCATION.ppsm__.pptx"
        if not os.path.exists(path):
            self.skipTest("sample file not available")
        from pptx import Presentation
        data = open(path, "rb").read()
        normalized, _ = _normalize_pptx_bytes(data)
        prs = Presentation(io.BytesIO(normalized))
        self.assertGreaterEqual(len(prs.slides), 1)


class TestExtractOleNativeInfo(unittest.TestCase):
    """Tests for _extract_ole_native_info helper."""

    def test_extracts_filename_src_temp_path(self):
        """_extract_ole_native_info returns filename, src_path, temp_path."""
        stream = _build_ole10native_stream(
            filename="invoice.exe",
            src_path=r"C:\Users\attacker\Desktop\invoice.exe",
            temp_path=r"C:\Users\ATTACK~1\AppData\Local\Temp\invoice.exe",
        )
        mock_ole = _mock_olefile(stream)
        with patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            result = _extract_ole_native_info(b"fake_cfb")

        self.assertIsNotNone(result)
        self.assertEqual(result["filename"], "invoice.exe")
        self.assertEqual(result["src_path"], r"C:\Users\attacker\Desktop\invoice.exe")
        self.assertEqual(result["temp_path"], r"C:\Users\ATTACK~1\AppData\Local\Temp\invoice.exe")

    def test_is_link_false_for_embedded(self):
        """_extract_ole_native_info sets is_link=False when payload is present."""
        stream = _build_ole10native_stream("doc.exe", r"C:\x\doc.exe", r"C:\tmp\doc.exe")
        mock_ole = _mock_olefile(stream)
        with patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            result = _extract_ole_native_info(b"fake_cfb")

        self.assertFalse(result["is_link"])

    def test_actual_size_matches_payload(self):
        """_extract_ole_native_info surfaces actual_size of the embedded payload."""
        payload = b"MZ" + b"\x00" * 510
        stream = _build_ole10native_stream("evil.exe", r"C:\x\evil.exe", r"C:\tmp\evil.exe", payload=payload)
        mock_ole = _mock_olefile(stream)
        with patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            result = _extract_ole_native_info(b"fake_cfb")

        self.assertEqual(result["actual_size"], len(payload))

    def test_returns_none_when_no_ole10native_stream(self):
        """_extract_ole_native_info returns None when no Ole10Native stream exists."""
        mock_ole = MagicMock()
        mock_ole.__enter__ = lambda s: s
        mock_ole.__exit__ = MagicMock(return_value=False)
        mock_ole.exists.return_value = False
        with patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            result = _extract_ole_native_info(b"fake_cfb")

        self.assertIsNone(result)

    def test_returns_none_on_parse_error(self):
        """_extract_ole_native_info returns None gracefully on malformed data."""
        mock_ole = MagicMock()
        mock_ole.__enter__ = lambda s: s
        mock_ole.__exit__ = MagicMock(return_value=False)
        mock_ole.exists.side_effect = Exception("corrupt")
        with patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            result = _extract_ole_native_info(b"fake_cfb")

        self.assertIsNone(result)

    def test_ole_object_metadata_includes_native(self):
        """OLE object dicts in scanner event include 'native' key for Package objects."""
        stream = _build_ole10native_stream(
            "payload.exe", r"C:\x\payload.exe", r"C:\tmp\payload.exe"
        )
        mock_ole = _mock_olefile(stream)

        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
        shape.click_action.hyperlink.address = "http://example.com"
        buf = io.BytesIO()
        prs.save(buf)
        pptx_data = buf.getvalue()

        from strelka.scanners.scan_pptx import ScanPptx
        scanner = ScanPptx.__new__(ScanPptx)
        scanner.event = {}
        scanner.flags = []
        scanner.files = []
        scanner.coordinator = MagicMock()
        scanner.name = "ScanPptx"
        scanner.scanner_timeout = 30

        # Inject a fake OLE object with native info into ole_objects post-scan
        # by patching _extract_shape_ole_metadata to return a Package object
        fake_ole_meta = {
            "slide": 1, "shape": "Object 1", "prog_id": "Package",
            "is_activex": False, "control_type": None,
            "show_as_icon": True, "blob_size": 512,
            "blob": b"fake_cfb",
        }
        with patch("strelka.scanners.scan_pptx._extract_shape_ole_metadata", return_value=fake_ole_meta), \
             patch("strelka.scanners.scan_pptx.olefile.OleFileIO", return_value=mock_ole):
            scanner.scan(pptx_data, MagicMock(), {}, 9999999999)

        ole_objects = scanner.event.get("ole_objects", [])
        package_objs = [o for o in ole_objects if o.get("prog_id") == "Package"]
        self.assertTrue(len(package_objs) > 0, "Expected a Package OLE object")
        pkg = package_objs[0]
        self.assertIn("native", pkg)
        self.assertEqual(pkg["native"]["filename"], "payload.exe")


if __name__ == "__main__":
    unittest.main()
