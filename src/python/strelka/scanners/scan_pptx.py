"""
Enhanced ScanPptx Scanner for Strelka

This enhanced version adds PowerPoint-specific active content detection:
- Hover actions (in addition to click actions)
- ppaction:// URI parsing (program, macro, ole verbs)
- Text-run hyperlinks
- Presentation-level relationships (remote templates, frames)
- OLE/ActiveX object metadata with type classification
- URL type classification
- Detection flags for suspicious content

Does NOT duplicate:
- ZIP member extraction (handled by ScanZip)
- OLE payload parsing (handled by ScanOle)
"""

import io
import zipfile
from urllib.parse import urlparse, parse_qs

from lxml import etree
from pptx import Presentation

from strelka import strelka


# OOXML namespaces for relationship and action parsing
NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}
RID_ATTR = "{%s}id" % NS["r"]

# High-risk ActiveX/OLE ProgIDs
HIGH_RISK_PROGIDS = {
    "Shell.Explorer",      # Web browser control
    "WScript.Shell",       # Shell execution
    "Shell.Application",   # Shell application
    "WScript.Network",     # Network access
    "Scripting.FileSystemObject",  # File system access
}

# Known CLSIDs and related vulnerabilities
# Source: https://github.com/decalage2/oletools/blob/master/oletools/common/clsid.py
# License: BSD
KNOWN_CLSIDS = {
    '00000300-0000-0000-C000-000000000046': 'StdOleLink (embedded OLE object - Known Related to CVE-2017-0199, CVE-2017-8570, CVE-2017-8759 or CVE-2018-8174)',
    '00000303-0000-0000-C000-000000000046': 'File Moniker (may trigger CVE-2017-0199 or CVE-2017-8570)',
    '00000304-0000-0000-C000-000000000046': 'Item Moniker',
    '00000305-0000-0000-C000-000000000046': 'Anti Moniker',
    '00000306-0000-0000-C000-000000000046': 'Pointer Moniker',
    '00000308-0000-0000-C000-000000000046': 'Packager Moniker',
    '00000309-0000-0000-C000-000000000046': 'Composite Moniker (may trigger CVE-2017-8570)',
    '0000031A-0000-0000-C000-000000000046': 'Class Moniker',
    '00000535-0000-0010-8000-00AA006D2EA4': 'ADODB.RecordSet (may trigger CVE-2015-0097)',
    '00000FE0-8804-4CA8-8868-36F59DEFD14D': 'ZED! encrypted container',
    '0002034C-0000-0000-C000-000000000046': 'OutlookAttachMoniker',
    '0002034E-0000-0000-C000-000000000046': 'OutlookMessageMoniker',
    '00020810-0000-0000-C000-000000000046': 'Microsoft Excel.Sheet.5',
    '00020811-0000-0000-C000-000000000046': 'Microsoft Excel.Chart.5',
    '00020820-0000-0000-C000-000000000046': 'Microsoft Microsoft Excel 97-2003 Worksheet (Excel.Sheet.8)',
    '00020821-0000-0000-C000-000000000046': 'Microsoft Excel.Chart.8',
    '00020830-0000-0000-C000-000000000046': 'Microsoft Excel.Sheet.12',
    '00020832-0000-0000-C000-000000000046': 'Microsoft Excel sheet with macro enabled (Excel.SheetMacroEnabled.12)',
    '00020833-0000-0000-C000-000000000046': 'Microsoft Excel binary sheet with macro enabled (Excel.SheetBinaryMacroEnabled.12)',
    '00020900-0000-0000-C000-000000000046': 'Microsoft Word 6.0-7.0 Document (Word.Document.6)',
    '00020906-0000-0000-C000-000000000046': 'Microsoft Word 97-2003 Document (Word.Document.8)',
    '00020907-0000-0000-C000-000000000046': 'Microsoft Word Picture (Word.Picture.8)',
    '00020C01-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '00021401-0000-0000-C000-000000000046': 'Windows LNK Shortcut file',
    '00021700-0000-0000-C000-000000000046': 'Microsoft Equation 2.0 (Known Related to CVE-2017-11882 or CVE-2018-0802)',
    '00022601-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '00022602-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '00022603-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '0002CE02-0000-0000-C000-000000000046': 'Microsoft Equation 3.0 (Known Related to CVE-2017-11882 or CVE-2018-0802)',
    '0002CE03-0000-0000-C000-000000000046': 'MathType Equation Object',
    '0003000A-0000-0000-C000-000000000046': 'Bitmap Image',
    '0003000B-0000-0000-C000-000000000046': 'Microsoft Equation (Known Related to CVE-2017-11882 or CVE-2018-0802)',
    '0003000C-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '0003000D-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '0003000E-0000-0000-C000-000000000046': 'OLE Package Object (may contain and run any file)',
    '0004A6B0-0000-0000-C000-000000000046': 'Microsoft Equation 2.0 (Known Related to CVE-2017-11882 or CVE-2018-0802)',
    '000C1082-0000-0000-C000-000000000046': 'MSI Transform (mst)',
    '000C1084-0000-0000-C000-000000000046': 'MSI Windows Installer Package (msi)',
    '000C1086-0000-0000-C000-000000000046': 'MSI Patch Package (psp)',
    '048EB43E-2059-422F-95E0-557DA96038AF': 'Microsoft Powerpoint.Slide.12',
    '05741520-C4EB-440A-AC3F-9643BBC9F847': 'otkloadr.WRLoader (can be used to bypass ASLR after triggering an exploit)',
    '06290BD2-48AA-11D2-8432-006008C3FBFC': 'Factory bindable using IPersistMoniker (scripletfile)',
    '06290BD3-48AA-11D2-8432-006008C3FBFC': 'Script Moniker, aka Moniker to a Windows Script Component (may trigger CVE-2017-0199)',
    '0CF774D0-F077-11D1-B1BC-00C04F86C324': 'scrrun.dll - HTML File Host Encode Object (ProgID: HTML.HostEncode)',
    '0D43FE01-F093-11CF-8940-00A0C9054228': 'scrrun.dll - FileSystem Object (ProgID: Scripting.FileSystemObject)',
    '0E59F1D5-1FBE-11D0-8FF2-00A0D10038BC': 'MSScriptControl.ScriptControl (may trigger CVE-2015-0097)',
    '1461A561-24E8-4BA3-8D4A-FFEEF980556B': 'BCSAddin.Connect (potential exploit CVE-2016-0042 / MS16-014)',
    '14CE31DC-ABC2-484C-B061-CF3416AED8FF': 'Loads WUAEXT.DLL (Known Related to CVE-2015-6128)',
    '18A06B6B-2F3F-4E2B-A611-52BE631B2D22': 'Word.DocumentMacroEnabled.12 (DOCM)',
    '1D8A9B47-3A28-4CE2-8A4B-BD34E45BCEEB': 'UPnP.DescriptionDocument',
    '1EFB6596-857C-11D1-B16A-00C0F0283628': 'MSCOMCTL.TabStrip (may trigger CVE-2012-1856, CVE-2013-3906 - often used for heap spray)',
    '233C1507-6A77-46A4-9443-F871F945D258': 'Shockwave Control Objects',
    '23CE100B-1390-49D6-BA00-F17D3AEE149C': 'UmOutlookAddin.UmEvmCtrl (potential exploit document CVE-2016-0042 / MS16-014)',
    '29131539-2EED-1069-BF5D-00DD011186B7': 'IBM/Lotus Notes COM interface provided by NLSXBE.DLL (related to CVE-2021-27058)',
    '3018609E-CDBC-47E8-A255-809D46BAA319': 'SSCE DropTable Listener Object (can be used to bypass ASLR after triggering an exploit)',
    '3050F4D8-98B5-11CF-BB82-00AA00BDCE0B': 'HTML Application (may trigger CVE-2017-0199)',
    '33BD73C2-7BB4-48F4-8DBC-82B8B313AE16': 'osf.SandboxManager (Known Related To CVE-2015-1770)',
    '33FD0563-D81A-4393-83CC-0195B1DA2F91': 'UPnP.DescriptionDocumentEx',
    '394C052E-B830-11D0-9A86-00C04FD8DBF7': 'Loads ELSEXT.DLL (Known Related to CVE-2015-6128)',
    '3BA59FA5-41BF-4820-98E4-04645A806698': 'osf.SandboxContent (Known Related To CVE-2015-1770)',
    '41B9BE05-B3AF-460C-BF0B-2CDD44A093B1': 'Search.XmlContentFilter',
    '4315D437-5B8C-11D0-BD3B-00A0C911CE86': 'Device Moniker (Known Related to CVE-2016-0015)',
    '44F9A03B-A3EC-4F3B-9364-08E0007F21DF': 'Control.TaskSymbol (Known Related to CVE-2015-1642 & CVE-2015-2424)',
    '46E31370-3F7A-11CE-BED6-00AA00611080': 'Forms.MultiPage',
    '4C599241-6926-101B-9992-00000B65C6F9': 'Forms.Image (may trigger CVE-2015-2424)',
    '4D3263E4-CAB7-11D2-802A-0080C703929C': 'AutoCAD 2000-2002 Document',
    '5E4405B0-5374-11CE-8E71-0020AF04B1D7': 'AutoCAD R14 Document',
    '64818D10-4F9B-11CF-86EA-00AA00B929E8': 'Microsoft Powerpoint.Show.8',
    '64818D11-4F9B-11CF-86EA-00AA00B929E8': 'Microsoft Powerpoint.Slide.8',
    '66833FE6-8583-11D1-B16A-00C0F0283628': 'MSCOMCTL.Toolbar (Known Related to CVE-2012-0158 & CVE-2012-1856)',
    '6A221957-2D85-42A7-8E19-BE33950D1DEB': 'AutoCAD 2013 Document',
    '6AD4AE40-2FF1-4D88-B27A-F76FC7B40440': 'BCSAddin.ManageSolutionHelper (potential exploit CVE-2016-0042 / MS16-014)',
    '6E182020-F460-11CE-9BCD-00AA00608E01': 'Forms.Frame',
    '799ED9EA-FB5E-11D1-B7D6-00C04FC2AAE2': 'Microsoft.VbaAddin (Known Related to CVE-2016-0042)',
    '79EAC9D0-BAF9-11CE-8C82-00AA004BA90B': 'StdHlink',
    '79EAC9D1-BAF9-11CE-8C82-00AA004BA90B': 'StdHlinkBrowseContext',
    '79EAC9E0-BAF9-11CE-8C82-00AA004BA90B': 'URL Moniker (may trigger CVE-2017-0199, CVE-2017-8570, or CVE-2018-8174)',
    '79EAC9E2-BAF9-11CE-8C82-00AA004BA90B': '(http:) Asychronous Pluggable Protocol Handler',
    '79EAC9E3-BAF9-11CE-8C82-00AA004BA90B': '(ftp:) Asychronous Pluggable Protocol Handler',
    '79EAC9E5-BAF9-11CE-8C82-00AA004BA90B': '(https:) Asychronous Pluggable Protocol Handler',
    '79EAC9E6-BAF9-11CE-8C82-00AA004BA90B': '(mk:) Asychronous Pluggable Protocol Handler',
    '79EAC9E7-BAF9-11CE-8C82-00AA004BA90B': '(file:, local:) Asychronous Pluggable Protocol Handler',
    '7AABBB95-79BE-4C0F-8024-EB6AF271231C': 'AutoCAD 2007-2009 Document',
    '85131630-480C-11D2-B1F9-00C04F86C324': 'scrrun.dll - JS File Host Encode Object (ProgID: JSFile.HostEncode)',
    '85131631-480C-11D2-B1F9-00C04F86C324': 'scrrun.dll - VBS File Host Encode Object (ProgID: VBSFile.HostEncode)',
    '8627E73B-B5AA-4643-A3B0-570EDA17E3E7': 'UmOutlookAddin.ButtonBar (potential exploit document CVE-2016-0042 / MS16-014)',
    '88D969E5-F192-11D4-A65F-0040963251E5': 'Msxml2.DOMDocument.5.0',
    '88D969E6-F192-11D4-A65F-0040963251E5': 'Msxml2.FreeThreadedDOMDocument.5.0',
    '88D969E7-F192-11D4-A65F-0040963251E5': 'Msxml2.XMLSchemaCache.5.0',
    '88D969E8-F192-11D4-A65F-0040963251E5': 'Msxml2.XSLTemplate.5.0',
    '88D969E9-F192-11D4-A65F-0040963251E5': 'Msxml2.DSOControl.5.0',
    '88D969EA-F192-11D4-A65F-0040963251E5': 'Msxml2.XMLHTTP.5.0',
    '88D969EB-F192-11D4-A65F-0040963251E5': 'Msxml2.ServerXMLHTTP.5.0',
    '88D969EC-8B8B-4C3D-859E-AF6CD158BE0F': 'Msxml2.SAXXMLReader.5.0',
    '88D969EE-F192-11D4-A65F-0040963251E5': 'Msxml2.SAXAttributes.5.0',
    '88D969EF-F192-11D4-A65F-0040963251E5': 'Msxml2.MXXMLWriter.5.0',
    '88D969F0-F192-11D4-A65F-0040963251E5': 'Msxml2.MXHTMLWriter.5.0',
    '88D969F1-F192-11D4-A65F-0040963251E5': 'Msxml2.MXNamespaceManager.5.0',
    '88D969F5-F192-11D4-A65F-0040963251E5': 'Msxml2.MXDigitalSignature.5.0',
    '88D96A0C-F192-11D4-A65F-0040963251E5': 'SAX XML Reader 6.0 (msxml6.dll)',
    '8E75D913-3D21-11D2-85C4-080009A0C626': 'AutoCAD 2004-2006 Document',
    '9181DC5F-E07D-418A-ACA6-8EEA1ECB8E9E': 'MSCOMCTL.TreeCtrl (may trigger CVE-2012-0158)',
    '975797FC-4E2A-11D0-B702-00C04FD8DBF7': 'Loads ELSEXT.DLL (Known Related to CVE-2015-6128)',
    '978C9E23-D4B0-11CE-BF2D-00AA003F40D0': 'Microsoft Forms 2.0 Label (Forms.Label.1)',
    '996BF5E0-8044-4650-ADEB-0B013914E99C': 'MSCOMCTL.ListViewCtrl (may trigger CVE-2012-0158)',
    '9C38ED61-D565-4728-AEEE-C80952F0ECDE': 'Virtual Disk Service Loader - vdsldr.exe (related to MS Office click-to-run issue CVE-2021-27058)',
    'A08A033D-1A75-4AB6-A166-EAD02F547959': 'otkloadr WRAssembly Object (can be used to bypass ASLR after triggering an exploit)',
    'B54F3741-5B07-11CF-A4B0-00AA004A55E8': 'vbscript.dll - VB Script Language (ProgID: VBS, VBScript)',
    'B801CA65-A1FC-11D0-85AD-444553540000': 'Adobe Acrobat Document - PDF file',
    'BDD1F04B-858B-11D1-B16A-00C0F0283628': 'MSCOMCTL.ListViewCtrl (may trigger CVE-2012-0158)',
    'C08AFD90-F2A1-11D1-8455-00A0C91F3880': 'ShellBrowserWindow',
    'C62A69F0-16DC-11CE-9E98-00AA00574A4F': 'Forms.Form',
    'C74190B6-8589-11D1-B16A-00C0F0283628': 'MSCOMCTL.TreeCtrl (may trigger CVE-2012-0158)',
    'CCD068CD-1260-4AEA-B040-A87974EB3AEF': 'UmOutlookAddin.RoomsCTP (potential exploit document CVE-2016-0042 / MS16-014)',
    'CDDBCC7C-BE18-4A58-9CBF-D62A012272CE': 'osf.Sandbox (Known Related To CVE-2015-1770)',
    'CDF1C8AA-2D25-43C7-8AFE-01F73A3C66DA': 'UmOutlookAddin.InspectorContext (potential exploit document CVE-2016-0042 / MS16-014)',
    'CF4F55F4-8F87-4D47-80BB-5808164BB3F8': 'Microsoft Powerpoint.Show.12',
    'D27CDB6E-AE6D-11CF-96B8-444553540000': 'Shockwave Flash Object (may trigger many CVEs)',
    'D27CDB70-AE6D-11CF-96B8-444553540000': 'Shockwave Flash Object (may trigger many CVEs)',
    'D50FED35-0A08-4B17-B3E0-A8DD0EDE375D': 'UmOutlookAddin.PlayOnPhoneDlg (potential exploit document CVE-2016-0042 / MS16-014)',
    'D7053240-CE69-11CD-A777-00DD01143C57': 'Microsoft Forms 2.0 CommandButton',
    'D70E31AD-2614-49F2-B0FC-ACA781D81F3E': 'AutoCAD 2010-2012 Document',
    'D93CE8B5-3BF8-462C-A03F-DED2730078BA': 'Loads WUAEXT.DLL (Known Related to CVE-2015-6128)',
    'DD9DA666-8594-11D1-B16A-00C0F0283628': 'MSCOMCTL.ImageComboCtrl (may trigger CVE-2014-1761)',
    'DF630910-1C1D-11D0-AE36-8C0F5E000000': 'pythoncomloader27.dll (related to CVE-2021-27058)',
    'DFEAF541-F3E1-4C24-ACAC-99C30715084A': 'Silverlight Objects',
    'E5CA59F5-57C4-4DD8-9BD6-1DEEEDD27AF4': 'InkEd.InkEdit',
    'E8CC4CBE-FDFF-11D0-B865-00A0C9081C1D': 'MSDAORA.1',
    'E8CC4CBF-FDFF-11D0-B865-00A0C9081C1D': 'Loads OCI.DLL (Known Related to CVE-2015-6128)',
    'ECABAFC6-7F19-11D2-978E-0000F8757E2A': 'New Moniker',
    'ECABAFC9-7F19-11D2-978E-0000F8757E2A': 'Loads MQRT.DLL (Known Related to CVE-2015-6128)',
    'ECABB0C7-7F19-11D2-978E-0000F8757E2A': 'SOAP Moniker (may trigger CVE-2017-8759)',
    'ECF44975-786E-462F-B02A-CBCCB1A2C4A2': 'UmOutlookAddin.FormRegionContext (potential exploit CVE-2016-0042 / MS16-014)',
    'F20DA720-C02F-11CE-927B-0800095AE340': 'OLE Package Object (may contain and run any file)',
    'F414C260-6AC0-11CF-B6D1-00AA00BBBB58': 'jscript.dll - JScript Language (ProgID: ECMAScript, JavaScript, JScript, LiveScript)',
    'F4754C9B-64F5-4B40-8AF4-679732AC0607': 'Microsoft Word Document (Word.Document.12)',
    'F959DBBB-3867-41F2-8E5F-3B8BEFAA81B3': 'UmOutlookAddin.FormRegionAddin (potential exploit document CVE-2016-0042 / MS16-014)',
    '00024512-0000-0000-C000-000000000046': 'RefEdit.Ctrl',
    '0002E500-0000-0000-C000-000000000046': 'Microsoft Office Chart 9.0',
    '0002E510-0000-0000-C000-000000000046': 'Microsoft Office Spreadsheet 9.0',
    '0002E520-0000-0000-C000-000000000046': 'Microsoft Office PivotTable 9.0',
    '0002E530-0000-0000-C000-000000000046': 'Microsoft Office Data Source Control 9.0',
    '0002E531-0000-0000-C000-000000000046': 'Microsoft Office Record Navigation Control 9.0',
    '0002E532-0000-0000-C000-000000000046': 'Microsoft Office Expand Control 9.0',
    '0002E551-0000-0000-C000-000000000046': 'Microsoft Office Spreadsheet 10.0',
    '0002E552-0000-0000-C000-000000000046': 'Microsoft Office PivotTable 10.0',
    '0002E553-0000-0000-C000-000000000046': 'Microsoft Office Data Source Control 10.0',
    '0002E554-0000-0000-C000-000000000046': 'Microsoft Office Record Navigation Control 10.0',
    '0002E556-0000-0000-C000-000000000046': 'Microsoft Office Chart 10.0',
    '0006F023-0000-0000-C000-000000000046': 'Microsoft Outlook Recipient Control',
    '0006F024-0000-0000-C000-000000000046': 'Microsoft Outlook Rich Format Control',
    '0006F063-0000-0000-C000-000000000046': 'Microsoft Outlook View Control',
    '009541A0-3B81-101C-92F3-040224009C02': 'Kodak Image Admin Control',
    '02BF25D5-8C17-4B23-BC80-D3488ABDDC6B': 'QuickTime Object',
    '04A1E553-FE36-4FDE-865E-344194E69424': 'Microsoft InkPicture Control',
    '0713E8A2-850A-101B-AFC0-4210102A8DA7': 'Microsoft TreeView Control, version 5.0 (SP2)',
    '0713E8D2-850A-101B-AFC0-4210102A8DA7': 'Microsoft ProgressBar Control, version 5.0 (SP2)',
    '0ECD9B64-23AA-11D0-B351-00A0C9055D8E': 'Microsoft Hierarchical FlexGrid Control 6.0 (SP4)',
    '166B1BCA-3F9C-11CF-8075-444553540000': 'Shockwave ActiveX Control',
    '20DD1B9E-87C4-11D1-8BE3-0000F8754DA1': 'Microsoft Date and Time Picker Control 6.0 (SP4)',
    '232E456A-87C3-11D1-8BE3-0000F8754DA1': 'Microsoft MonthView Control 6.0 (SP4)',
    '2C247F23-8591-11D1-B16A-00C0F0283628': 'Microsoft ImageList Control 6.0 (SP4)',
    '35053A22-8589-11D1-B16A-00C0F0283628': 'Microsoft ProgressBar Control 6.0 (SP4)',
    '373FF7F0-EB8B-11CD-8820-08002B2F4F5A': 'Microsoft Slider Control, version 5.0 (SP2)',
    '3A2B370C-BA0A-11D1-B137-0000F8753F5D': 'Microsoft Chart Control, version 6.0 (SP3)',
    '3B7C8860-D78F-101B-B9B5-04021C009402': 'Microsoft Rich Textbox Control 6.0 (SP6)',
    '48E59293-9880-11CF-9754-00AA00C00908': 'Microsoft Internet Transfer Control 6.0 (SP4)',
    '58DA8D8A-9D6A-101B-AFC0-4210102A8DA7': 'Microsoft ListView Control, version 5.0 (SP2)',
    '58DA8D8F-9D6A-101B-AFC0-4210102A8DA7': 'Microsoft ImageList Control, version 5.0 (SP2)',
    '603C7E80-87C2-11D1-8BE3-0000F8754DA1': 'Microsoft UpDown Control 6.0 (SP4)',
    '612A8624-0FB3-11CE-8747-524153480004': 'Microsoft Toolbar Control, version 5.0 (SP2)',
    '6262D3A0-531B-11CF-91F6-C2863C385E30': 'Microsoft FlexGrid Control, version 6.0',
    '6B7E638F-850A-101B-AFC0-4210102A8DA7': 'Microsoft StatusBar Control, version 5.0 (SP2)',
    '6BF52A52-394A-11D3-B153-00C04F79FAA6': 'Windows Media Player',
    '6D940280-9F11-11CE-83FD-02608C3EC08A': 'Kodak Image Edit Control',
    '6D940285-9F11-11CE-83FD-02608C3EC08A': 'Kodak Image Annotation Control',
    '79176FB0-B7F2-11CE-97EF-00AA006D2776': 'Microsoft Forms 2.0 SpinButton',
    '84926CA0-2941-101C-816F-0E6013114B7F': 'Kodak Image Scan Control',
    '8856F961-340A-11D0-A96B-00C04FD705A2': 'Microsoft Web Browser',
    '8BD21D10-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 TextBox',
    '8BD21D20-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 ListBox',
    '8BD21D30-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 ComboBox',
    '8BD21D40-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 CheckBox',
    '8BD21D50-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 OptionButton',
    '8BD21D60-EC42-11CE-9E0D-00AA006002F3': 'Microsoft Forms 2.0 ToggleButton',
    '8E3867A3-8586-11D1-B16A-00C0F0283628': 'Microsoft StatusBar Control 6.0 (SP4)',
    '9ED94440-E5E8-101B-B9B5-444553540000': 'Microsoft TabStrip Control, version 5.0 (SP2)',
    'B09DE715-87C1-11D1-8BE3-0000F8754DA1': 'Microsoft Animation Control 6.0 (SP4)',
    'BDC217C5-ED16-11CD-956C-0000C04E4C0A': 'Microsoft Tabbed Dialog Control 6.0 (SP5)',
    'C4D2D8E0-D1DD-11CE-940F-008029004347': 'System Monitor Control',
    'CA8A9780-280D-11CF-A24D-444553540000': 'Adobe Acrobat Control for ActiveX',
    'CFCDAA03-8BE4-11CF-B84B-0020AFBBCCFA': 'RealPlayer G2 Control',
    'D45FD31B-5C6E-11D1-9EC1-00C04FD7081F': 'Microsoft Agent Control 2.0',
    'DFD181E0-5E2F-11CE-A449-00AA004A803D': 'Microsoft Forms 2.0 ScrollBar',
    'E1A6B8A0-3603-101C-AC6E-040224009C02': 'Kodak Image Thumbnail Control',
    'EAE50EB0-4A62-11CE-BED6-00AA00611080': 'Microsoft Forms 2.0 TabStrip',
    'F08DF954-8592-11D1-B16A-00C0F0283628': 'Microsoft Slider Control 6.0 (SP4)',
    'F9043C85-F6F2-101A-A3C9-08002B2F49FB': 'Microsoft Common Dialog Control, version 6.0',
    'FE38753A-44A3-11D1-B5B7-0000C09000C4': 'Microsoft Flat Scrollbar Control 6.0 (SP4)',
}


def _lookup_clsid(clsid):
    """
    Look up CLSID description and check if dangerous.

    Args:
        clsid: CLSID string (with or without braces)

    Returns:
        tuple of (description, is_dangerous)
    """
    # Normalize CLSID format (remove braces, convert to uppercase)
    normalized = clsid.strip().strip('{}').upper()

    desc = KNOWN_CLSIDS.get(normalized, None)

    # Check if dangerous based on description keywords
    is_dangerous = False
    if desc:
        dangerous_indicators = [
            "CVE-", "may trigger", "Known Related", "potential exploit",
            "may contain and run", "OLE Package Object"
        ]
        is_dangerous = any(indicator in desc for indicator in dangerous_indicators)

    return (desc, is_dangerous)


def _extract_activex_controls_from_zip(pptx_io):
    """
    Extract ActiveX controls from ppt/activeX/ directory in PPTX file.

    ActiveX controls are stored in ppt/activeX/ and referenced via
    'control' relationship type from slides.

    Args:
        pptx_io: BytesIO object of the PPTX file (positioned at start)

    Returns:
        list of ActiveX control dictionaries with CLSID lookups
    """
    activex_controls = []

    try:
        with zipfile.ZipFile(pptx_io, 'r') as zf:
            # Find all activeX XML files
            activex_files = [name for name in zf.namelist() if 'activeX/activeX' in name and name.endswith('.xml')]

            for ax_file in activex_files:
                try:
                    ax_xml = zf.read(ax_file)
                    root = etree.fromstring(ax_xml)

                    # Namespace for ActiveX
                    ax_ns = {"ax": "http://schemas.microsoft.com/office/2006/activeX"}

                    classid = root.get("{http://schemas.microsoft.com/office/2006/activeX}classid", "")
                    persistence = root.get("{http://schemas.microsoft.com/office/2006/activeX}persistence", "")

                    # Extract properties
                    properties = {}
                    for prop in root.findall(".//ax:ocxPr", namespaces=ax_ns):
                        name = prop.get("{http://schemas.microsoft.com/office/2006/activeX}name")
                        value = prop.get("{http://schemas.microsoft.com/office/2006/activeX}value")
                        if name and value:
                            properties[name] = value

                    # Look up CLSID in known database
                    clsid_desc, is_dangerous = _lookup_clsid(classid)

                    activex_controls.append({
                        "type": "activex_control",
                        "classid": classid,
                        "clsid_desc": clsid_desc,
                        "persistence": persistence,
                        "properties": properties,
                        "is_activex": True,
                        "is_high_risk": is_dangerous,
                        "source_file": ax_file,
                    })
                except Exception:
                    # If we can't parse, still record the file
                    activex_controls.append({
                        "type": "activex_control",
                        "source_file": ax_file,
                        "is_activex": True,
                        "parse_error": True,
                    })
    except Exception:
        pass

    return activex_controls


def _parse_ppaction(action_url):
    """
    Parse ppaction:// URLs to extract verb and query parameters.

    Args:
        action_url: String like "ppaction://program?file=malware.exe"

    Returns:
        dict with 'ppaction_url', 'verb', and 'fields'
    """
    if not action_url or not action_url.startswith("ppaction://"):
        return {"ppaction_url": None, "verb": None, "fields": {}}

    parsed = urlparse(action_url)
    return {
        "ppaction_url": action_url,
        "verb": (parsed.netloc or None),
        "fields": {
            k: (v[0] if len(v) == 1 else v)
            for k, v in parse_qs(parsed.query).items()
        },
    }


def _get_relationship_info(part, rid):
    """
    Get relationship target URL and external status.

    Args:
        part: The presentation part containing relationships
        rid: Relationship ID (e.g., "rId5")

    Returns:
        tuple: (target_url, is_external)
    """
    if not rid:
        return None, None

    try:
        rel = part.rels.get(rid)
        if not rel:
            return None, None

        target = getattr(rel, "target_ref", None) or getattr(rel, "_target", None)
        is_external = getattr(rel, "is_external", False)

        return (str(target) if target else None), bool(is_external)
    except Exception:
        return None, None


def _extract_shape_actions(shape, slide_num):
    """
    Extract both click and hover actions from a shape.

    Args:
        shape: python-pptx shape object
        slide_num: Slide number for context

    Returns:
        list of action dictionaries
    """
    actions = []
    part = shape.part
    shape_name = getattr(shape, "name", "Unknown")

    # XPath queries for click and hover actions
    for trigger, xpath in (
        ("click", ".//a:hlinkClick"),
        ("hover", ".//a:hlinkHover"),
    ):
        try:
            for node in etree.ElementBase.xpath(shape.element, xpath, namespaces=NS):
                rid = node.get(RID_ATTR)
                action_url = node.get("action")

                # Parse ppaction if present
                pp_info = _parse_ppaction(action_url)
                target, is_external = _get_relationship_info(part, rid)

                actions.append({
                    "slide": slide_num,
                    "shape": shape_name,
                    "trigger": trigger,
                    "verb": pp_info["verb"],
                    "ppaction_url": pp_info["ppaction_url"],
                    "fields": pp_info["fields"],
                    "rid": rid or None,
                    "target": target,
                    "is_external": is_external,
                })
        except Exception:
            # Skip malformed action elements
            pass

    return actions


def _extract_text_run_hyperlinks(shape, slide_num):
    """
    Extract hyperlinks from text runs within a shape.

    Args:
        shape: python-pptx shape object
        slide_num: Slide number for context

    Returns:
        list of relationship dictionaries
    """
    hyperlinks = []
    shape_name = getattr(shape, "name", "Unknown")

    if not shape.has_text_frame:
        return hyperlinks

    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if hasattr(run, "hyperlink") and run.hyperlink:
                    address = getattr(run.hyperlink, "address", None)
                    if address:
                        hyperlinks.append({
                            "type": "hyperlink",
                            "target": address,
                            "is_external": True,
                            "location": f"slide_{slide_num}_shape_{shape_name}",
                        })
    except Exception:
        pass

    return hyperlinks


def _extract_shape_ole_metadata(shape, slide_num):
    """
    Extract OLE/ActiveX object metadata without parsing payload.

    This reports metadata for detection purposes. Actual OLE content
    extraction is handled by ScanZip → ScanOle.

    Args:
        shape: python-pptx shape object
        slide_num: Slide number for context

    Returns:
        dict with OLE/ActiveX metadata or None
    """
    try:
        ole = shape.ole_format
    except (AttributeError, KeyError):
        return None

    if not ole:
        return None

    shape_name = getattr(shape, "name", "Unknown")
    prog_id = getattr(ole, "prog_id", None)
    blob = getattr(ole, "blob", None)

    # Detect ActiveX controls
    is_activex = False
    control_type = None
    is_high_risk = False

    if prog_id:
        # Check for ActiveX prefixes
        activex_prefixes = [
            "Forms.",
            "MSForms.",
            "MSComCtl",
            "Shell.Explorer",
            "WScript.",
            "Scripting.",
        ]
        is_activex = any(prog_id.startswith(prefix) for prefix in activex_prefixes)

        # Extract control type from ProgID
        if is_activex and "." in prog_id:
            parts = prog_id.split(".")
            if len(parts) >= 2:
                control_type = parts[1]  # e.g., "CommandButton" from "Forms.CommandButton.1"

        # Check for high-risk ProgIDs
        is_high_risk = any(
            prog_id.startswith(risk_id) for risk_id in HIGH_RISK_PROGIDS
        )

    return {
        "slide": slide_num,
        "shape": shape_name,
        "prog_id": prog_id,
        "is_activex": is_activex,
        "control_type": control_type,
        "is_high_risk": is_high_risk,
        "show_as_icon": getattr(ole, "show_as_icon", None),
        "blob_size": len(blob) if blob else None,
    }


def _classify_relationship_type(rel_type):
    """
    Classify relationship type for security analysis.

    Args:
        rel_type: Relationship type URI

    Returns:
        string classification
    """
    type_map = {
        "attachedTemplate": "remote_template",
        "image": "external_image",
        "video": "external_video",
        "audio": "external_audio",
        "frame": "remote_frame",
        "hyperlink": "hyperlink",
        "oleObject": "ole_object",
        "control": "activex_control",
    }

    if not rel_type:
        return "unknown"

    # Extract the last part of the relationship type URI
    rel_name = rel_type.split("/")[-1] if "/" in rel_type else rel_type

    return type_map.get(rel_name, "external_content")


def _extract_presentation_relationships(presentation):
    """
    Extract presentation-level relationships (templates, external content).

    Args:
        presentation: python-pptx Presentation object

    Returns:
        list of relationship dictionaries
    """
    relationships = []

    try:
        for rid, rel in presentation.part.rels.items():
            rel_type = getattr(rel, "reltype", "")
            target = getattr(rel, "target_ref", None) or getattr(rel, "_target", None)
            is_external = getattr(rel, "is_external", False)

            # Check for specific relationship types that indicate security concerns
            # Remote templates and external links should always be reported
            is_remote_template = "attachedTemplate" in rel_type or "externalLink" in rel_type

            # Report external relationships or remote templates
            if (is_external or is_remote_template) and target:
                relationships.append({
                    "type": _classify_relationship_type(rel_type),
                    "target": str(target),
                    "is_external": is_external or is_remote_template,
                    "rid": rid,
                    "location": "presentation",
                })
    except Exception:
        pass

    return relationships


def _extract_slide_relationships(slide, slide_num):
    """
    Extract slide-level external relationships (images, videos, frames).

    Args:
        slide: python-pptx slide object
        slide_num: Slide number for context

    Returns:
        list of relationship dictionaries
    """
    relationships = []

    try:
        for rid, rel in slide.part.rels.items():
            # Only report external relationships
            if getattr(rel, "is_external", False):
                target = getattr(rel, "target_ref", None) or getattr(rel, "_target", None)
                rel_type = getattr(rel, "reltype", None)

                if target:
                    relationships.append({
                        "type": _classify_relationship_type(rel_type),
                        "target": str(target),
                        "is_external": True,
                        "rid": rid,
                        "location": f"slide_{slide_num}",
                    })
    except Exception:
        pass

    return relationships


class ScanPptx(strelka.Scanner):
    """
    Collects metadata, extracts text, and detects active content from PPTX files.

    This scanner focuses on PowerPoint-specific security indicators:
    - Click and hover actions (including ppaction:// URIs)
    - Hyperlinks and external relationships
    - OLE/ActiveX object metadata with type classification
    - Remote templates and external content

    File extraction is handled by ScanZip.
    OLE content parsing is handled by ScanOle.

    Options:
        extract_text: Boolean that determines if document text should be
            extracted as a child file.
            Defaults to False.
    """

    def scan(self, data, file, options, expire_at):
        extract_text = options.get("extract_text", False)

        with io.BytesIO(data) as pptx_io:
            try:
                pptx_doc = Presentation(pptx_io)

                # Extract core properties
                self.event.update({
                    "author": pptx_doc.core_properties.author,
                    "category": pptx_doc.core_properties.category,
                    "comments": pptx_doc.core_properties.comments,
                    "content_status": pptx_doc.core_properties.content_status,
                    "created": (
                        int(pptx_doc.core_properties.created.strftime("%s"))
                        if pptx_doc.core_properties.created is not None
                        else None
                    ),
                    "identifier": pptx_doc.core_properties.identifier,
                    "keywords": pptx_doc.core_properties.keywords,
                    "language": pptx_doc.core_properties.language,
                    "last_modified_by": pptx_doc.core_properties.last_modified_by,
                    "last_printed": (
                        int(pptx_doc.core_properties.last_printed.strftime("%s"))
                        if pptx_doc.core_properties.last_printed is not None
                        else None
                    ),
                    "modified": (
                        int(pptx_doc.core_properties.modified.strftime("%s"))
                        if pptx_doc.core_properties.modified is not None
                        else None
                    ),
                    "revision": pptx_doc.core_properties.revision,
                    "subject": pptx_doc.core_properties.subject,
                    "title": pptx_doc.core_properties.title,
                    "version": pptx_doc.core_properties.version,
                    "slide_count": len(pptx_doc.slides),
                    "word_count": 0,
                    "image_count": 0,
                })

                # Initialize collection lists
                extracted_text = [] if extract_text else None
                extracted_notes = []
                extracted_urls = []  # Backward compatibility
                all_relationships = []
                all_actions = []
                all_ole_objects = []

                # Extract presentation-level relationships
                pres_rels = _extract_presentation_relationships(pptx_doc)
                all_relationships.extend(pres_rels)

                # Extract ActiveX controls from ppt/activeX/ directory
                try:
                    pptx_io.seek(0)  # Reset to beginning for ZIP parsing
                    activex_controls = _extract_activex_controls_from_zip(pptx_io)
                    if activex_controls:
                        all_ole_objects.extend(activex_controls)
                except Exception:
                    self.flags.append("activex_extraction_error")

                # Process each slide
                for slide_num, slide in enumerate(pptx_doc.slides, start=1):
                    # Extract slide-level external relationships
                    try:
                        slide_rels = _extract_slide_relationships(slide, slide_num)
                        all_relationships.extend(slide_rels)
                    except Exception:
                        self.flags.append(f"slide_{slide_num}_rel_error")

                    # Process shapes
                    for shape in slide.shapes:
                        try:
                            # Count images
                            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                                self.event["image_count"] += 1
                        except Exception:
                            pass

                        # Process text frames
                        try:
                            if shape.has_text_frame:
                                for para in shape.text_frame.paragraphs:
                                    # Collect text for extraction
                                    if extract_text and para.text:
                                        extracted_text.append(para.text)

                                    # Count words
                                    for run in para.runs:
                                        text = run.text.strip()
                                        if text:
                                            self.event["word_count"] += len(text.split())
                        except Exception:
                            pass

                        # Extract actions (click and hover)
                        try:
                            shape_actions = _extract_shape_actions(shape, slide_num)
                            all_actions.extend(shape_actions)
                        except Exception:
                            pass

                        # Extract text-run hyperlinks
                        try:
                            text_hyperlinks = _extract_text_run_hyperlinks(shape, slide_num)
                            all_relationships.extend(text_hyperlinks)
                        except Exception:
                            pass

                        # Extract shape click action URL (backward compatibility)
                        try:
                            if hasattr(shape, "click_action") and shape.click_action:
                                if (
                                    shape.click_action.hyperlink
                                    and shape.click_action.hyperlink.address
                                ):
                                    url = shape.click_action.hyperlink.address
                                    extracted_urls.append(url)
                        except Exception:
                            pass

                        # Extract OLE/ActiveX object metadata
                        try:
                            ole_metadata = _extract_shape_ole_metadata(shape, slide_num)
                            if ole_metadata:
                                all_ole_objects.append(ole_metadata)
                        except Exception:
                            pass

                    # Extract speaker notes
                    try:
                        if slide.has_notes_slide:
                            notes_text = slide.notes_slide.notes_text_frame.text.strip()
                            if notes_text:
                                extracted_notes.append(notes_text)
                    except Exception:
                        pass

                # Add URLs to event (backward compatibility)
                if extracted_urls:
                    self.event["urls"] = extracted_urls

                # Add enhanced relationship data
                if all_relationships:
                    # Deduplicate by (type, target, location)
                    seen = set()
                    unique_rels = []
                    for rel in all_relationships:
                        key = (rel["type"], rel["target"], rel.get("location", ""))
                        if key not in seen:
                            seen.add(key)
                            unique_rels.append(rel)

                    self.event["relationships"] = unique_rels

                # Add action data
                if all_actions:
                    # Deduplicate by (slide, shape, trigger, verb, target)
                    seen = set()
                    unique_actions = []
                    for action in all_actions:
                        key = (
                            action["slide"],
                            action["shape"],
                            action["trigger"],
                            action["verb"],
                            action["target"],
                        )
                        if key not in seen:
                            seen.add(key)
                            unique_actions.append(action)

                    self.event["actions"] = unique_actions

                # Add OLE/ActiveX metadata
                if all_ole_objects:
                    self.event["ole_objects"] = all_ole_objects

                # Add detection flags
                self.event["has_hover_actions"] = any(
                    a["trigger"] == "hover" for a in all_actions
                )
                self.event["has_ppaction_program"] = any(
                    a["verb"] == "program" for a in all_actions
                )
                self.event["has_ppaction_macro"] = any(
                    a["verb"] == "macro" for a in all_actions
                )
                self.event["has_ppaction_ole"] = any(
                    a["verb"] == "ole" for a in all_actions
                )
                self.event["has_external_relationships"] = any(
                    r.get("is_external", False) for r in all_relationships
                )
                self.event["has_remote_template"] = any(
                    r["type"] == "remote_template" for r in all_relationships
                )
                self.event["has_activex_controls"] = any(
                    obj.get("is_activex", False) for obj in all_ole_objects
                )
                self.event["has_high_risk_ole"] = any(
                    obj.get("is_high_risk", False) for obj in all_ole_objects
                )

                # Add notes array
                if extracted_notes:
                    self.event["notes"] = extracted_notes

                # Upload extracted text as single batch
                if extract_text and extracted_text:
                    extract_file = strelka.File(
                        name="text",
                        source=self.name,
                    )

                    text_content = "\n".join(extracted_text)
                    for c in strelka.chunk_string(text_content):
                        self.upload_to_coordinator(
                            extract_file.pointer,
                            c,
                            expire_at,
                        )

                    self.files.append(extract_file)

            except ValueError:
                self.flags.append("value_error")
            except zipfile.BadZipFile:
                self.flags.append("bad_zip")
            except Exception as e:
                self.flags.append("processing_error")
